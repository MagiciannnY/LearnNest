# app/services/file_parser.py

import os
from typing import Literal
from pathlib import Path

from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation

def extract_text_from_file(file_path: Path, file_type: Literal["txt", "pdf", "docx", "pptx", "md"]) -> str:
    match file_type:
        case "txt" | "md":
            return file_path.read_text(encoding="utf-8")
        case "pdf":
            return extract_pdf_text(file_path)
        case "docx":
            return "\n".join(para.text for para in Document(file_path).paragraphs)
        case "pptx":
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        case _:
            raise ValueError(f"Unsupported file type: {file_type}")

